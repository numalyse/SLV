import sys
import json
import os
from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN, MSO_AUTO_SIZE
from pptx.dml.color import RGBColor

def format_time(time_val):
    """
    Formate le temps si le C++ envoie des millisecondes brutes, 
    ou renvoie le texte direct si le C++ a déjà formaté la chaîne.
    """
    if isinstance(time_val, str):
        return time_val
    
    seconds = time_val / 1000.0
    h = int(seconds // 3600)
    m = int((seconds % 3600) // 60)
    s = int(seconds % 60)
    f = int((seconds - int(seconds)) * 100)
    return f"{h:02d}:{m:02d}:{s:02d}:{f:02d}"

def main():
    if len(sys.argv) < 2:
        print("Erreur: Chemin du fichier JSON manquant.")
        sys.exit(1)

    json_path = sys.argv[1]
    
    with open(json_path, 'r', encoding='utf-8') as f:
        data = json.load(f)

    temp_dir = data['tempDir']
    dst_path = data['dstPath']
    if not dst_path.endswith('.pptx'):
        dst_path += '.pptx'

    shots = data['shots']
    total_shots = len(shots)

    try:
        prs = Presentation()
        prs.slide_width = Inches(13.33)
        prs.slide_height = Inches(7.5)
        
        blank_slide_layout = prs.slide_layouts[6]
        slide_width = prs.slide_width
        slide_height = prs.slide_height

        slide = prs.slides.add_slide(blank_slide_layout)

        textbox_width = slide_width - Inches(1)
        textbox_left = (slide_width - textbox_width) / 2
        
        title_height = Inches(2)
        title_box = slide.shapes.add_textbox(textbox_left, (slide_height - title_height) / 2, textbox_width, title_height)
        title_tf = title_box.text_frame
        title_tf.vertical_anchor = MSO_ANCHOR.MIDDLE

        title_p = title_tf.paragraphs[0]
        title_p.alignment = PP_ALIGN.CENTER
        title_p.font.bold = True
        title_p.font.size = Pt(32)
        title_p.font.color.rgb = RGBColor.from_string("EB4034") # Rouge
        title_p.text = "Étude cinématographique"

        # Slides des plans
        plan_box_top = Inches(0.5)
        plan_box_height = Inches(0.5)
        note_box_height = Inches(1)
        offset_between_boxes = Inches(0.2)

        for idx, shot in enumerate(shots):
            slide = prs.slides.add_slide(blank_slide_layout)

            y_offset = plan_box_top + plan_box_height + offset_between_boxes

            # Titre
            plan_box = slide.shapes.add_textbox(textbox_left, plan_box_top, textbox_width, plan_box_height)
            plan_tf = plan_box.text_frame
            plan_tf.vertical_anchor = MSO_ANCHOR.MIDDLE

            plan_p = plan_tf.paragraphs[0]
            plan_p.font.bold = True
            plan_p.font.color.rgb = RGBColor.from_string("1F497D")
            

            time_str = format_time(shot.get('start', 0))
            end_str = format_time(shot.get('duration', 0))
            plan_title = shot.get('title', f"Plan {idx+1}")
            
            plan_p.text = f"[Plan {idx+1}] {plan_title} - Début : {time_str} / Durée : {end_str}"

            img_path = os.path.join(temp_dir, shot['image'])
            if os.path.exists(img_path):
                try:
                    image = Image.open(img_path)
                    img_w, img_h = image.size
                    dpi = image.info.get('dpi', (96, 96))[0]
                    
                    img_width_emu = Inches(img_w / dpi)
                    img_height_emu = Inches(img_h / dpi)

                    max_w = slide_width - Inches(1)
                    max_h = slide_height - plan_box_height - note_box_height - y_offset - offset_between_boxes

                    scale = min(max_w / img_width_emu, max_h / img_height_emu)
                    new_w = img_width_emu * scale
                    new_h = img_height_emu * scale

                    x_pos = (slide_width - new_w) / 2
                    y_pos = y_offset

                    slide.shapes.add_picture(img_path, x_pos, y_pos, width=new_w, height=new_h)
                    
                    y_offset += new_h + offset_between_boxes
                    
                except Exception as e:
                    print(f"Erreur lors du traitement de l'image {idx} : {e}", file=sys.stderr)


            note_text = shot.get('note', "").strip()
            
            note_box = slide.shapes.add_textbox(textbox_left, y_offset, textbox_width, note_box_height)
            note_box.height = note_box_height
            note_tf = note_box.text_frame
            note_tf.word_wrap = True

            note_p = note_tf.paragraphs[0]
            note_p.alignment = PP_ALIGN.LEFT
            note_p.font.size = Pt(16)

            if note_text:
                note_p.text = note_text
            else:
                note_tf.auto_size = MSO_AUTO_SIZE.NONE
                note_p.text = "\u00A0"
                note_box.height = note_box_height

            percent = int(((idx + 1) / total_shots) * 100)
            print(f"PROGRESS:{percent}", flush=True)


        prs.save(dst_path)
        print("Export PPTX terminé avec succès !", flush=True)

    except Exception as e:
        print(f"Erreur Fatale lors de l'exportation PPTX : {e}", file=sys.stderr)
        sys.exit(1)

if __name__ == "__main__":
    main()